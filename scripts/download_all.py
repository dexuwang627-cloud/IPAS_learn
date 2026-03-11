#!/usr/bin/env python3
"""
遞迴掃描 Google Drive IPAS 資料夾，下載所有 PDF/DOCX/PPTX 並提取文字。
用法: python scripts/download_all.py [--root-id FOLDER_ID] [--generate]
"""
import subprocess
import json
import os
import sys
import shutil
import argparse
import hashlib
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

GWS = os.path.expanduser("~/.cargo/bin/gws")
OUTPUT_DIR = "data/texts"
TEMP_DIR = "data/_temp"
ROOT_FOLDER_ID = "1QutCz68R8-7UX2tPD7AzRI5zxAg0msiG"

SUPPORTED_MIME = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
    "application/msword": "doc",
}
FOLDER_MIME = "application/vnd.google-apps.folder"


def list_folder(folder_id: str) -> list[dict]:
    """列出 Google Drive 資料夾內容"""
    try:
        r = subprocess.run(
            [GWS, "drive", "files", "list", "--params",
             json.dumps({
                 "q": f'"{folder_id}" in parents',
                 "fields": "files(id,name,mimeType)",
                 "pageSize": 200,
             })],
            capture_output=True, text=True, timeout=30
        )
        return json.loads(r.stdout).get("files", [])
    except Exception as e:
        print(f"    ⚠️ 無法列出資料夾 {folder_id}: {e}")
        return []


def scan_folder(folder_id: str, path: str = "") -> list[dict]:
    """遞迴掃描資料夾，回傳所有支援的檔案列表"""
    results = []
    files = list_folder(folder_id)
    for f in files:
        fp = f"{path}/{f['name']}" if path else f['name']
        if f["mimeType"] == FOLDER_MIME:
            results.extend(scan_folder(f["id"], fp))
        elif f["mimeType"] in SUPPORTED_MIME:
            results.append({
                "id": f["id"],
                "name": f["name"],
                "path": fp,
                "ext": SUPPORTED_MIME[f["mimeType"]],
                "chapter": _infer_chapter(fp),
            })
    return results


def _infer_chapter(path: str) -> str:
    """從路徑推斷章節名稱"""
    parts = path.split("/")
    # 找 L2xx 開頭的資料夾或檔名
    for p in parts:
        for prefix in ["L211", "L212", "L213", "L221", "L222", "L223"]:
            if prefix in p:
                return prefix
    # 用次深的資料夾名
    if len(parts) >= 2:
        return parts[-2]
    return parts[0].rsplit(".", 1)[0]


def download_file(file_id: str, dest_path: str) -> bool:
    """從 Google Drive 下載檔案"""
    os.makedirs(TEMP_DIR, exist_ok=True)
    # 清空 temp
    for f in os.listdir(TEMP_DIR):
        os.remove(os.path.join(TEMP_DIR, f))
    try:
        subprocess.run(
            [GWS, "drive", "files", "get", "--params",
             json.dumps({"fileId": file_id, "alt": "media"})],
            capture_output=True, text=True, timeout=120,
            cwd=TEMP_DIR,
        )
        for f in os.listdir(TEMP_DIR):
            if f.startswith("download"):
                shutil.move(os.path.join(TEMP_DIR, f), dest_path)
                return True
        return False
    except Exception as e:
        print(f"      下載失敗: {e}")
        return False


def extract_text_pdf(path: str) -> str:
    import pdfplumber
    try:
        with pdfplumber.open(path) as pdf:
            pages = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages)
    except Exception:
        return ""


def extract_text_docx(path: str) -> str:
    from docx import Document
    try:
        doc = Document(path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts)
    except Exception:
        return ""


def extract_text_pptx(path: str) -> str:
    from pptx import Presentation
    try:
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            parts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            t = cell.text.strip()
                            if t:
                                parts.append(t)
        return "\n".join(parts)
    except Exception:
        return ""


EXTRACTORS = {
    "pdf": extract_text_pdf,
    "docx": extract_text_docx,
    "pptx": extract_text_pptx,
    "doc": lambda p: "",  # .doc 較難處理，跳過
}


def file_hash(name: str, file_id: str) -> str:
    """產生唯一短 hash 避免檔名衝突"""
    return hashlib.md5(file_id.encode()).hexdigest()[:6]


def main():
    parser = argparse.ArgumentParser(description="遞迴下載 Google Drive IPAS 資料並提取文字")
    parser.add_argument("--root-id", default=ROOT_FOLDER_ID, help="根資料夾 ID")
    parser.add_argument("--generate", action="store_true", help="下載後自動產題")
    parser.add_argument("--num-choice", type=int, default=3)
    parser.add_argument("--num-tf", type=int, default=2)
    args = parser.parse_args()

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(TEMP_DIR, exist_ok=True)

    # 讀取已處理的檔案
    done_path = os.path.join(OUTPUT_DIR, "_processed.json")
    processed = {}
    if os.path.exists(done_path):
        with open(done_path, "r", encoding="utf-8") as f:
            processed = json.load(f)

    print("🔍 掃描 Google Drive IPAS 資料夾...")
    all_files = scan_folder(args.root_id, "IPAS")
    print(f"   找到 {len(all_files)} 個檔案")

    new_files = [f for f in all_files if f["id"] not in processed]
    print(f"   其中 {len(new_files)} 個尚未處理\n")

    if not new_files:
        print("✅ 全部檔案已處理完畢，無需下載")
        return

    extracted_count = 0
    for i, fi in enumerate(new_files, 1):
        print(f"[{i}/{len(new_files)}] 📥 {fi['name']} ({fi['ext']})")

        # 下載
        h = file_hash(fi["name"], fi["id"])
        temp_path = os.path.join(TEMP_DIR, f"file_{h}.{fi['ext']}")
        if not download_file(fi["id"], temp_path):
            processed[fi["id"]] = {"name": fi["name"], "status": "download_failed"}
            continue

        # 提取文字
        extractor = EXTRACTORS.get(fi["ext"])
        if not extractor:
            processed[fi["id"]] = {"name": fi["name"], "status": "unsupported"}
            continue

        text = extractor(temp_path)
        if os.path.exists(temp_path):
            os.remove(temp_path)

        if not text or len(text.strip()) < 50:
            print(f"   ⏭️  文字太少，跳過")
            processed[fi["id"]] = {"name": fi["name"], "status": "no_text"}
            continue

        # 儲存文字
        safe_name = fi["name"].rsplit(".", 1)[0]
        # 避免檔名衝突
        txt_name = f"{safe_name}_{h}.txt"
        out_path = os.path.join(OUTPUT_DIR, txt_name)
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(f"# Chapter: {fi['chapter']}\n\n{text}")

        processed[fi["id"]] = {
            "name": fi["name"],
            "txt": txt_name,
            "chapter": fi["chapter"],
            "chars": len(text),
            "status": "ok",
        }
        extracted_count += 1
        print(f"   ✅ {len(text)} 字 → {txt_name}")

    # 儲存處理紀錄
    with open(done_path, "w", encoding="utf-8") as f:
        json.dump(processed, f, ensure_ascii=False, indent=2)

    print(f"\n📊 完成！本次提取 {extracted_count} 個檔案的文字")
    total_ok = sum(1 for v in processed.values() if v.get("status") == "ok")
    print(f"   累計已處理: {total_ok} 個檔案有文字")

    # 清理 temp
    if os.path.exists(TEMP_DIR):
        shutil.rmtree(TEMP_DIR, ignore_errors=True)

    # 自動產題
    if args.generate and extracted_count > 0:
        print("\n🤖 開始自動產題...")
        from database import init_db, insert_question, get_stats
        from services.question_generator import generate_from_text_files

        init_db()
        questions = generate_from_text_files(
            texts_dir=OUTPUT_DIR,
            num_choice=args.num_choice,
            num_tf=args.num_tf,
        )
        inserted = 0
        for q in questions:
            try:
                insert_question(q)
                inserted += 1
            except Exception:
                pass
        stats = get_stats()
        print(f"\n✅ 產生 {len(questions)} 題，寫入 {inserted} 題")
        print(f"📊 題庫總計: {stats['total']} 題")


if __name__ == "__main__":
    main()
